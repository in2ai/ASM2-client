import csv
import io
import logging
from pathlib import Path

import chardet
import requests
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from dropbox.exceptions import ApiError
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader

from src.config.search_config import GRAPH
from src.utils.helpers import safe_execute


def extract_docx_text(data: bytes) -> str:
    """Extraer texto de documento Word (.docx)"""
    try:
        doc = DocxDocument(io.BytesIO(data))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                ]
                if row_text:
                    tables_text.append(" | ".join(row_text))
        return "\n".join(paragraphs + tables_text).strip()
    except Exception as e:
        logging.error(f"Error while extracting DOCX: {e}")
        return None


def extract_pptx_text(data: bytes) -> list[str]:
    """Extraer texto de presentación PowerPoint (.pptx), una cadena por diapositiva."""
    try:
        prs = Presentation(io.BytesIO(data))
        slides_text = []

        for slide in prs.slides:
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            slides_text.append("\n".join(slide_parts).strip())

        return slides_text
    except Exception as e:
        logging.error(f"Error while extracting PPTX: {e}")
        return None


def extract_excel_text(data: bytes) -> str:
    """Extraer texto de hoja de cálculo Excel (.xlsx, .xls)"""
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        text_rows = []
        for sheet in wb.worksheets:
            text_rows.append(f"Hoja: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(row_values).strip()
                if row_text and row_text != " | " * (len(row_values) - 1):
                    text_rows.append(row_text)
        return "\n".join(text_rows).strip()
    except Exception as e:
        logging.error(f"Error while extracting Excel: {e}")
        return None


def extract_html_text(data: bytes) -> str:
    """Extraer texto de documento HTML"""
    try:
        soup = BeautifulSoup(data, "lxml")
        # Eliminar elementos script y style
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator="\n")
        # Limpiar espacios en blanco
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines).strip()
    except Exception as e:
        logging.error(f"Error while extracting HTML: {e}")
        return None


def extract_csv_text(data: bytes) -> str:
    """Extraer texto de archivo CSV"""
    try:
        # Detectar encoding automáticamente (soporta UTF-8, Latin-1, Windows-1252, etc.)
        detected = chardet.detect(data)
        encoding = detected.get("encoding") or "utf-8"
        text = data.decode(encoding, errors="ignore")

        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            row_text = " | ".join([cell.strip() for cell in row if cell.strip()])
            if row_text:
                rows.append(row_text)
        return "\n".join(rows).strip()
    except Exception as e:
        logging.error(f"Error while extracting CSV: {e}")
        return None


def extract_pdf_text(data: bytes) -> list[str]:
    """Extraer texto de documento PDF, una cadena por página."""
    try:
        fh = io.BytesIO(data)
        reader = PdfReader(fh)
        return [(p.extract_text() or "").strip() for p in reader.pages]
    except Exception as e:
        logging.error(f"Error while extracting PDF: {e}")
        return None


class VDBFile:
    def __init__(self, metadata):
        self.metadata = metadata

    def download(self, path: str) -> Path: ...
    def get_text(self) -> str | list[str]: ...


class GoogleDriveFile(VDBFile):
    GOOGLE_EXPORT_MIME = {
        "application/vnd.google-apps.document": "text/plain",
        "application/vnd.google-apps.spreadsheet": "text/csv",
        "application/vnd.google-apps.presentation": "text/plain",
    }

    GOOGLE_EXPORTS = {
        "application/vnd.google-apps.document": ("text/plain", ".txt"),
        "application/vnd.google-apps.spreadsheet": ("text/csv", ".csv"),
        "application/vnd.google-apps.presentation": ("text/plain", ".txt"),
    }

    DIRECT_DOWNLOAD_MIME = {
        "application/pdf",
        "text/plain",
        "text/markdown",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "text/html",
        "text/csv",
    }

    def __init__(self, metadata, service):
        super().__init__(metadata)
        self.service = service

    def get_data(self):
        file_id = self.metadata["id"]
        mime_type = self.metadata["mimeType"]

        if mime_type in self.GOOGLE_EXPORT_MIME:
            export_mime = self.GOOGLE_EXPORT_MIME[mime_type]
            return safe_execute(
                self.service.files().export(fileId=file_id, mimeType=export_mime)
            )

        return safe_execute(self.service.files().get_media(fileId=file_id))

    def get_text(self) -> str | list[str] | None:
        mime_type = self.metadata["mimeType"]
        data = self.get_data()

        if data is None:
            return None

        if mime_type == "application/pdf":
            return extract_pdf_text(data)

        if mime_type == "application/vnd.google-apps.spreadsheet":
            return extract_csv_text(data)

        if mime_type in {
            "application/vnd.google-apps.document",
            "application/vnd.google-apps.presentation",
            "text/plain",
            "text/markdown",
        }:
            return data.decode("utf-8", errors="ignore")

        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_docx_text(data)

        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_pptx_text(data)

        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return extract_excel_text(data)

        if mime_type == "text/html":
            return extract_html_text(data)

        if mime_type == "text/csv":
            return extract_csv_text(data)

        return None

    def _get_download_spec(self):
        file_id = self.metadata["id"]
        mime_type = self.metadata["mimeType"]

        if mime_type in self.GOOGLE_EXPORTS:
            export_mime, extension = self.GOOGLE_EXPORTS[mime_type]
            request = self.service.files().export(fileId=file_id, mimeType=export_mime)
            return request, extension

        name = self.metadata.get("name", file_id)
        extension = Path(name).suffix
        request = self.service.files().get_media(fileId=file_id)
        return request, extension

    def download(self, directory: str | Path) -> Path:
        directory = Path(directory)
        directory.mkdir(parents=True, exist_ok=True)

        request, extension = self._get_download_spec()
        file_id = self.metadata["id"]
        path = directory / f"{file_id}{extension}"

        data = safe_execute(request)
        if data is None:
            raise ValueError("No data available for download")

        path.write_bytes(data)
        return path


class DropboxFile(VDBFile):
    def __init__(self, metadata, service):
        super().__init__(metadata)
        self.service = service

    def get_text(self) -> str | list[str]:
        file_id = self.metadata["id"]
        path_lower = self.metadata["path_lower"]
        mime_type = self.metadata["mimeType"]

        try:
            meta, resp = self.service.files_download(file_id or path_lower)
            data = resp.content

            # Archivos PDF
            if mime_type == "application/pdf":
                return extract_pdf_text(data)

            # Texto plano y Markdown
            elif mime_type in ("text/plain", "text/markdown"):
                return data.decode("utf-8", errors="ignore")

            # Documentos Word
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                return extract_docx_text(data)

            # Presentaciones PowerPoint
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):
                return extract_pptx_text(data)

            # Hojas de cálculo Excel
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ):
                return extract_excel_text(data)

            # Archivos HTML
            elif mime_type == "text/html":
                return extract_html_text(data)

            # Archivos CSV
            elif mime_type == "text/csv":
                return extract_csv_text(data)

        except ApiError:
            return None

        return None


def _ms_headers(token_dict):
    return {"Authorization": f"Bearer {token_dict['access_token']}"}


class OnedriveFile(VDBFile):
    def __init__(self, metadata, token):
        super().__init__(metadata)
        self.token = token

    def get_text(self) -> str | list[str]:
        item_id = self.metadata["id"]
        mime = self.metadata.get("mimeType", "").lower()

        headers = _ms_headers(self.token)
        url = f"{GRAPH}/me/drive/items/{item_id}/content"
        r = requests.get(url, headers=headers, timeout=60)

        if r.status_code == 403:
            return None

        r.raise_for_status()
        data = r.content

        # Archivos PDF
        if mime == "application/pdf":
            return extract_pdf_text(data)

        # Documentos Word
        elif (
            mime
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            return extract_docx_text(data)

        # Presentaciones PowerPoint
        elif (
            mime
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            return extract_pptx_text(data)

        # Hojas de cálculo Excel
        elif (
            mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            return extract_excel_text(data)

        # Archivos HTML
        elif mime == "text/html":
            return extract_html_text(data)

        # Archivos CSV
        elif mime == "text/csv":
            return extract_csv_text(data)

        # Texto plano y Markdown
        elif mime in ("text/plain", "text/markdown"):
            return data.decode("utf-8", errors="ignore")

        # Tipo no soportado
        return None
